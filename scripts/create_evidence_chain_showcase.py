from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = "docs/ppt/liangda-evidence-chain-showcase-editable.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def rgb(hex_color):
    h = hex_color.replace("#", "")
    return RGBColor.from_string(h)


C = {
    "bg": rgb("F6F8FC"), "ink": rgb("18253D"), "muted": rgb("65738B"),
    "navy": rgb("1237B8"), "blue": rgb("2F6BFF"), "blue_pale": rgb("EAF1FF"),
    "teal": rgb("00A99A"), "teal_pale": rgb("E8F8F5"),
    "orange": rgb("F59E0B"), "orange_pale": rgb("FFF4DD"),
    "red": rgb("D85B62"), "red_pale": rgb("FFF0F1"),
    "line": rgb("D8E1EF"), "white": rgb("FFFFFF"), "soft": rgb("EEF3FA"),
}


def shape(kind, x, y, w, h, fill=None, line=None, radius=False):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.8)
    return sh


def text(value, x, y, w, h, size=12, color=None, bold=False, align=PP_ALIGN.LEFT,
         font="Microsoft YaHei", valign=MSO_ANCHOR.TOP, margin=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin); tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value; run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color or C["ink"]
    return box


slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C["bg"]

# Header
text("04 / 核心功能展示", 0.62, 0.35, 3.2, 0.35, 23, C["navy"], True)
text("证据链可解释 AI", 0.62, 0.83, 5.4, 0.48, 31, C["ink"], True)
text("每一次生成、推荐与拦截，都能回溯到真实依据", 0.64, 1.35, 6.5, 0.24, 11, C["muted"])
text("LIANGDA HEALTH  /  TRUST LAYER", 10.2, 0.5, 2.45, 0.2, 8, C["navy"], True, PP_ALIGN.RIGHT)

# Main left product mockup panel
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 1.86, 7.55, 4.75, C["white"], C["line"])
text("真实交互：证据随结果生成，解释在对话中呈现", 0.9, 2.08, 6.7, 0.25, 13, C["ink"], True)
text("Evidence appears with the answer, not after the fact", 0.9, 2.39, 5.8, 0.2, 8.5, C["muted"])

# Browser frame
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.78, 6.98, 3.36, C["soft"], C["line"])
shape(MSO_SHAPE.RECTANGLE, 0.9, 2.78, 6.98, 0.34, rgb("E9EFF8"), rgb("E9EFF8"))
for i, col in enumerate(("E85D6A", "F5B942", "16A085")):
    shape(MSO_SHAPE.OVAL, 1.08 + i * 0.17, 2.9, 0.09, 0.09, rgb(col), rgb(col))
text("chat / family-health-assistant", 1.65, 2.89, 3.2, 0.12, 7.4, C["muted"])

# Chat area
shape(MSO_SHAPE.RECTANGLE, 1.08, 3.34, 3.48, 2.58, C["white"], C["line"])
text("家庭健康助手", 1.3, 3.55, 2.2, 0.2, 10.5, C["ink"], True)
text("针对李阿姨近期血糖波动，今晚建议\n减少精制主食，优先选择高纤维搭配。", 1.3, 3.95, 2.9, 0.56, 10, C["ink"])
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 1.3, 4.78, 2.82, 0.45, C["blue_pale"], C["blue"])
text("查看本次回复的证据链", 1.48, 4.91, 2.45, 0.15, 8.7, C["blue"], True)
text("生成依据   推荐依据", 1.3, 5.49, 2.7, 0.18, 8.3, C["muted"])

# Evidence panel in mockup
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 4.72, 3.34, 2.98, 2.58, C["white"], C["line"])
text("证据链", 4.98, 3.55, 1.0, 0.22, 11, C["ink"], True)
text("生成依据", 6.23, 3.57, 1.0, 0.16, 8.3, C["blue"], True, PP_ALIGN.RIGHT)
shape(MSO_SHAPE.RECTANGLE, 4.98, 3.92, 2.46, 0.52, C["blue_pale"], C["blue"])
text("报告事实 · 5 月体检报告 p3\n空腹血糖 6.8 mmol/L", 5.14, 4.03, 2.1, 0.25, 8.2, C["ink"])
shape(MSO_SHAPE.RECTANGLE, 4.98, 4.58, 2.46, 0.52, C["teal_pale"], C["teal"])
text("互动记忆 · 家庭成员偏好\n晚餐倾向清淡少油", 5.14, 4.69, 2.1, 0.25, 8.2, C["ink"])
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 4.98, 5.33, 2.46, 0.35, C["orange_pale"], C["orange"])
text("来源真实 · 可追溯 · 可核验", 5.12, 5.43, 2.1, 0.12, 7.6, C["orange"], True)

# Callout markers over the product frame
for n, x, y, col in (("A", 4.22, 4.83, C["blue"]), ("B", 7.46, 3.63, C["teal"]), ("C", 7.44, 5.52, C["orange"])):
    shape(MSO_SHAPE.OVAL, x, y, 0.28, 0.28, col, col)
    text(n, x, y + 0.055, 0.28, 0.12, 8, C["white"], True, PP_ALIGN.CENTER)

# Right rail cards
text("三层解释能力", 8.55, 1.98, 2.1, 0.25, 15, C["ink"], True)
text("让 AI 的结果可解释、可验证、可控", 8.55, 2.29, 3.9, 0.18, 9.5, C["muted"])

cards = [
    ("01", "生成依据", "为什么这样建议？", "报告事实  ·  健康画像  ·  家庭记忆", C["blue"], C["blue_pale"]),
    ("02", "推荐依据", "为什么推荐这款？", "健康标签  ·  营养成分  ·  适配人群", C["teal"], C["teal_pale"]),
    ("03", "安全拦截", "为什么没有继续回答？", "过敏禁忌  ·  特殊人群  ·  风险规则", C["red"], C["red_pale"]),
]
for i, (num, title, question, detail, accent, pale) in enumerate(cards):
    y = 2.78 + i * 1.12
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8.52, y, 4.28, 0.92, C["white"], C["line"])
    shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8.72, y + 0.19, 0.42, 0.42, pale, pale)
    text(num, 8.72, y + 0.325, 0.42, 0.12, 8.5, accent, True, PP_ALIGN.CENTER)
    text(title, 9.34, y + 0.14, 1.3, 0.18, 12, C["ink"], True)
    text(question, 10.72, y + 0.15, 1.75, 0.16, 9, accent, True, PP_ALIGN.RIGHT)
    text(detail, 9.34, y + 0.48, 3.1, 0.15, 8.5, C["muted"])

# Right-side takeaway
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 8.52, 6.2, 4.28, 0.42, C["navy"], C["navy"])
text("不是“相信 AI”，而是“看见依据”", 8.78, 6.33, 3.75, 0.14, 9.4, C["white"], True, PP_ALIGN.CENTER)

# Footer pathway
shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 6.92, 12.18, 0.38, C["navy"], C["navy"])
text("真实数据", 0.94, 7.04, 0.68, 0.12, 8.7, C["white"], True)
text("→", 1.76, 7.01, 0.22, 0.14, 11, rgb("BFD3FF"), True, PP_ALIGN.CENTER)
text("生成建议", 2.18, 7.04, 0.72, 0.12, 8.7, C["white"], True)
text("→", 3.08, 7.01, 0.22, 0.14, 11, rgb("BFD3FF"), True, PP_ALIGN.CENTER)
text("推荐商品", 3.48, 7.04, 0.72, 0.12, 8.7, C["white"], True)
text("→", 4.38, 7.01, 0.22, 0.14, 11, rgb("BFD3FF"), True, PP_ALIGN.CENTER)
text("安全校验", 4.78, 7.04, 0.72, 0.12, 8.7, C["white"], True)
text("→", 5.68, 7.01, 0.22, 0.14, 11, rgb("BFD3FF"), True, PP_ALIGN.CENTER)
text("可解释输出", 6.08, 7.04, 0.9, 0.12, 8.7, C["white"], True)
text("可追溯 · 可核验 · 可控", 10.52, 7.04, 1.95, 0.12, 8.5, rgb("9FE7DB"), True, PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
