#!/usr/bin/env python3
"""
drawio -> 可编辑 PowerPoint 像素级转换器

用法:
  python3 drawio_to_pptx.py [input.drawio] [output.pptx]
  未指定时取本目录下的 context_engineering_arch.*
"""
from __future__ import annotations

import math
import re
import sys
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu, Inches
from lxml import etree

# ---------------------------------------------------------------------------
# 配置
# ---------------------------------------------------------------------------
HERE = Path(__file__).parent
DEFAULT_IN = HERE / "context_engineering_arch.drawio"
DEFAULT_OUT = HERE / "context_engineering_arch.pptx"

SCALE = 1.0  # drawio 单位 -> ppt 内部单位 (1:1, 因 1 pt ≈ drawio 1 单位)


def resolve_paths() -> tuple[Path, Path]:
    inp = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_IN
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else inp.with_suffix(".pptx")
    return inp, out


# ---------------------------------------------------------------------------
# 工具: 颜色 / 样式
# ---------------------------------------------------------------------------

def hex_to_rgb(h: str | None) -> RGBColor | None:
    if not h:
        return None
    h = h.lstrip("#").lower()
    if not re.fullmatch(r"[0-9a-f]{6}", h):
        return None
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def parse_style(style: str) -> dict[str, str]:
    out: dict[str, str] = {}
    if not style:
        return out
    tokens = [t.strip() for t in style.split(";") if t.strip()]
    for tok in tokens:
        if "=" in tok:
            k, v = tok.split("=", 1)
            out[k.strip()] = v.strip()
        else:
            # 标志 (无值): 如 text, swimlane, rhombus
            out[tok] = "1"
    return out


# ---------------------------------------------------------------------------
# 文本 (含 HTML 简单解析) -> TextLine 列表
# ---------------------------------------------------------------------------

def _decode_entities(s: str) -> str:
    return (
        s.replace("&lt;", "<")
        .replace("&gt;", ">")
        .replace("&amp;", "&")
        .replace("&quot;", '"')
        .replace("&#xa;", "\n")
    )


@dataclass
class TextLine:
    runs: list[dict] = field(default_factory=list)
    bold: bool = False
    size: int | None = None
    color: str | None = None


class TextParser:
    """把 drawio value 字符串解析为 TextLine 列表, 支持简单 HTML"""
    def __init__(self, raw: str, default_size: int | None = None,
                 default_bold: bool = False, default_color: str | None = None) -> None:
        self.lines: list[TextLine] = []
        self._cur = TextLine(bold=default_bold, size=default_size, color=default_color)
        self._ib: bool | None = None
        self._is: int | None = None
        self._ic: str | None = None
        if raw:
            self._parse(_decode_entities(raw))
        else:
            self.lines.append(self._cur)

    def _flush(self) -> None:
        self.lines.append(self._cur)
        self._cur = TextLine(bold=self._cur.bold, size=self._cur.size, color=self._cur.color)

    def _resolve(self) -> tuple[bool, int | None, str | None]:
        bold = self._ib if self._ib is not None else self._cur.bold
        size = self._is if self._is is not None else self._cur.size
        color = self._ic if self._ic is not None else self._cur.color
        return bold, size, color

    def _add(self, txt: str) -> None:
        if not txt:
            return
        # 支持 token 内部含 \n - 自动按行切分 (避免 run.text 含 \n 导致 pptx 渲染异常)
        parts = txt.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                # 在切分点插入新行
                self._flush()
            if not part:
                continue
            bold, size, color = self._resolve()
            self._cur.runs.append({"text": part, "bold": bold, "size": size, "color": color})

    def _newline(self) -> None:
        self._flush()

    def _handle_tag(self, tag: str) -> None:
        t = tag.lower().strip()
        if t in ("<br>", "<br/>", "<br />"):
            self._newline()
            return
        if t == "<b>":
            self._ib = True
            return
        if t == "</b>":
            self._ib = None
            return
        if t.startswith("<div"):
            self._newline()
            return
        if t == "</div>":
            self._newline()
            return
        if t.startswith("<font"):
            m = re.search(r"font-size:\s*(\d+)px", t)
            if m:
                self._is = int(m.group(1))
            mc = re.search(r"color:\s*(#[0-9a-fA-F]{3,6})", t)
            if mc:
                self._ic = mc.group(1)
            mb = re.search(r"font-weight:\s*bold", t)
            if mb:
                self._ib = True
            return
        if t == "</font>":
            self._is = None
            self._ic = None
            self._ib = None

    def _parse(self, html: str) -> None:
        # 把换行 token 化
        tokens = re.split(r"(<[^>]+>)", html)
        for tok in tokens:
            if not tok:
                continue
            if tok.startswith("<"):
                self._handle_tag(tok)
            else:
                self._add(tok)
        self._flush()

    def get_lines(self) -> list[TextLine]:
        # 去掉首尾空白行
        ls = self.lines[:]
        while ls and not ls[0].runs:
            ls.pop(0)
        while ls and not ls[-1].runs:
            ls.pop()
        return ls or [TextLine()]


# ---------------------------------------------------------------------------
# 解析 drawio
# ---------------------------------------------------------------------------

@dataclass
class Cell:
    id: str
    parent: str | None
    style: dict[str, str]
    raw_value: str
    is_vertex: bool
    is_edge: bool
    source: str | None
    target: str | None
    geom_x: float
    geom_y: float
    w: float
    h: float
    waypoints: list[tuple[float, float]] = field(default_factory=list)


def load_cells(path: Path) -> dict[str, Cell]:
    tree = ET.parse(path)
    root = tree.getroot()
    cells: dict[str, Cell] = {}
    for mx in root.iter("mxCell"):
        cid = mx.get("id") or ""
        if cid in ("0", "1"):
            continue
        style = parse_style(mx.get("style") or "")
        value = mx.get("value") or ""
        is_vertex = mx.get("vertex") == "1"
        is_edge = mx.get("edge") == "1"
        geom = mx.find("mxGeometry")
        x = y = w = h = 0.0
        waypoints: list[tuple[float, float]] = []
        if geom is not None:
            try:
                x = float(geom.get("x", 0) or 0)
                y = float(geom.get("y", 0) or 0)
                w = float(geom.get("width", 0) or 0)
                h = float(geom.get("height", 0) or 0)
            except Exception:
                pass
            for arr in geom.findall("Array"):
                for pt in arr.findall("mxPoint"):
                    try:
                        px = float(pt.get("x", 0) or 0)
                        py = float(pt.get("y", 0) or 0)
                        waypoints.append((px, py))
                    except Exception:
                        pass
        # 跳过 drawio 内部 cell (层)
        cells[cid] = Cell(
            id=cid,
            parent=mx.get("parent"),
            style=style,
            raw_value=value,
            is_vertex=is_vertex,
            is_edge=is_edge,
            source=mx.get("source"),
            target=mx.get("target"),
            geom_x=x, geom_y=y, w=w, h=h,
            waypoints=waypoints,
        )
    return cells


# ---------------------------------------------------------------------------
# 计算绝对坐标
# ---------------------------------------------------------------------------

def compute_absolute(cells: dict[str, Cell]) -> dict[str, tuple[float, float]]:
    abs_pos: dict[str, tuple[float, float]] = {}

    def resolve(cid: str, depth: int = 0) -> tuple[float, float]:
        if cid in abs_pos:
            return abs_pos[cid]
        if depth > 50:
            return (0.0, 0.0)
        c = cells.get(cid)
        if c is None:
            return (0.0, 0.0)
        if c.parent in ("0", "1", None) or c.parent not in cells:
            abs_pos[cid] = (c.geom_x, c.geom_y)
            return abs_pos[cid]
        px, py = resolve(c.parent, depth + 1)
        ax, ay = px + c.geom_x, py + c.geom_y
        abs_pos[cid] = (ax, ay)
        return abs_pos[cid]

    for cid in cells:
        resolve(cid)
    return abs_pos


# ---------------------------------------------------------------------------
# 形状分类
# ---------------------------------------------------------------------------

def shape_kind(c: Cell) -> str:
    """返回 pptx 形状种类"""
    s = c.style
    if "swimlane" in s:
        return "swimlane"
    if "text" in s:
        return "textonly"
    if s.get("shape") == "process":
        return "process"
    if s.get("rounded", "0") == "1":
        return "rounded"
    return "rect"


def alignment_of(c: Cell) -> PP_ALIGN:
    a = c.style.get("align", "center").lower()
    return {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}.get(a, PP_ALIGN.CENTER)


def valign_of(c: Cell) -> MSO_ANCHOR:
    a = c.style.get("verticalAlign", "middle").lower()
    return {"top": MSO_ANCHOR.TOP, "bottom": MSO_ANCHOR.BOTTOM, "middle": MSO_ANCHOR.MIDDLE}.get(a, MSO_ANCHOR.MIDDLE)


def font_size_of(c: Cell) -> int:
    try:
        return int(float(c.style.get("fontSize", "12")))
    except Exception:
        return 12


def is_bold_of(c: Cell) -> bool:
    return c.style.get("fontStyle", "0") == "1"


def font_color_of(c: Cell) -> str | None:
    return c.style.get("fontColor")


def is_edge_label(c: Cell) -> bool:
    """edgeLabel 是附着在 edge 上的虚拟顶点, 由我们自己在边绘制时处理"""
    return "edgelabel" in c.style


# ---------------------------------------------------------------------------
# OOXML 辅助
# ---------------------------------------------------------------------------

def set_fill(shape, rgb: RGBColor | None) -> None:
    if rgb is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb: RGBColor | None, width_pt: float = 1.0,
             dashed: bool = False, dash_pattern: str = "dash") -> None:
    """设置轮廓"""
    line = shape.line
    if rgb is None or width_pt <= 0:
        # 设为 0 宽度 + 无填充, 防止默认阴影
        try:
            line.fill.background()
        except Exception:
            pass
        line.width = Pt(0.5)
    else:
        line.color.rgb = rgb
        line.width = Pt(max(width_pt, 0.25))
    if dashed:
        try:
            ln = shape.line._get_or_add_ln()
            for child in list(ln):
                if etree.QName(child).localname == "prstDash":
                    ln.remove(child)
            pd = etree.SubElement(ln, qn("a:prstDash"))
            pd.set("val", dash_pattern if dash_pattern else "dash")
        except Exception:
            pass


def set_no_shadow(shape) -> None:
    """去掉形状默认阴影"""
    try:
        spPr = shape._element.spPr
        for child in list(spPr):
            if etree.QName(child).localname == "effectLst":
                spPr.remove(child)
        etree.SubElement(spPr, qn("a:effectLst"))
    except Exception:
        pass


def add_arrow_head_end(connector) -> None:
    """给 connector 末端加三角形箭头"""
    try:
        ln = connector.line._get_or_add_ln()
        for child in list(ln):
            tag = etree.QName(child).localname
            if tag == "tailEnd":
                ln.remove(child)
        te = etree.SubElement(ln, qn("a:tailEnd"))
        te.set("type", "triangle")
        te.set("w", "med")
        te.set("len", "med")
    except Exception:
        pass


def _line_pattern(style: dict[str, str]) -> tuple[bool, str]:
    if style.get("dashed") == "1":
        return True, style.get("dashPattern") or "dash"
    return False, "solid"


# ---------------------------------------------------------------------------
# 应用文本
# ---------------------------------------------------------------------------

def apply_text_to(shape, parser: TextParser, align: PP_ALIGN, vanchor: MSO_ANCHOR,
                  default_color: str | None = None) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = vanchor

    lines = parser.get_lines()
    if not lines:
        return

    fallback_rgb = hex_to_rgb(default_color) if default_color else RGBColor(0x0f, 0x17, 0x2a)

    p0 = tf.paragraphs[0]
    p0.alignment = align
    _runs_to_paragraph(p0, lines[0], default_color=fallback_rgb)

    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        _runs_to_paragraph(p, ln, default_color=fallback_rgb)


def _runs_to_paragraph(p, line: TextLine, default_color: RGBColor | None = None) -> None:
    # 清空默认 run
    for r in list(p.runs):
        r._r.getparent().remove(r._r)

    if not line.runs:
        run = p.add_run()
        run.text = ""
        f = run.font
        if line.size:
            f.size = Pt(line.size)
        if default_color is not None:
            f.color.rgb = default_color
        return

    for r in line.runs:
        run = p.add_run()
        run.text = r.get("text", "")
        f = run.font
        if r.get("size"):
            f.size = Pt(int(r["size"]))
        elif line.size:
            f.size = Pt(line.size)
        if r.get("bold") is True:
            f.bold = True
        elif r.get("bold") is False:
            f.bold = False
        # 显式颜色 - 防止 p:style 的 fontRef=lt1 让文本看不见
        col = r.get("color") or line.color
        if col:
            c = hex_to_rgb(col)
            if c is not None:
                f.color.rgb = c
        elif default_color is not None:
            f.color.rgb = default_color


# 给 TextParser 加 raw_has_color
TextParser.raw_has_color = False
_orig_init = TextParser.__init__


def _patched_init(self, raw, default_size=None, default_bold=False, default_color=None):
    _orig_init(self, raw, default_size, default_bold, default_color)
    # 检查 raw 中是否有 color=
    self.raw_has_color = bool(re.search(r"color\s*:", raw or "", flags=re.I))

TextParser.__init__ = _patched_init


# ---------------------------------------------------------------------------
# 创建形状
# ---------------------------------------------------------------------------

def add_shape_node(slide, kind: str, x_pt: float, y_pt: float, w_pt: float, h_pt: float):
    if kind == "swimlane":
        return slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))
    if kind == "rounded":
        return slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))
    if kind == "process":
        # drawio 中 shape=process 通常是带轻微圆角的矩形, 而非流程图图标
        return slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))
    if kind == "textonly":
        return slide.shapes.add_textbox(Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))


def style_node(slide_shape, c: Cell) -> None:
    s = c.style
    fc = hex_to_rgb(s.get("fillColor"))
    sc = hex_to_rgb(s.get("strokeColor"))
    width = float(s.get("strokeWidth", "1") or 1)
    dashed, pat = _line_pattern(s)
    set_fill(slide_shape, fc)
    set_line(slide_shape, sc, width_pt=width, dashed=dashed, dash_pattern=pat)
    set_no_shadow(slide_shape)


def text_to_node(slide_shape, c: Cell) -> None:
    parser = TextParser(c.raw_value,
                        default_size=font_size_of(c),
                        default_bold=is_bold_of(c),
                        default_color=font_color_of(c))
    apply_text_to(slide_shape, parser, alignment_of(c), valign_of(c))


def create_vertex(slide, c: Cell, abs_x: float, abs_y: float) -> None:
    kind = shape_kind(c)
    shp = add_shape_node(slide, kind,
                         (abs_x) * SCALE,
                         (abs_y) * SCALE,
                         c.w * SCALE,
                         c.h * SCALE)
    style_node(shp, c)
    text_to_node(shp, c)


def create_swimlane(slide, c: Cell, abs_x: float, abs_y: float) -> None:
    """swimlane: 容器本体 + 独立标题条 (避免和正文内容重叠)"""
    s = c.style
    start_size = float(s.get("startSize", "32") or 32)
    # 主体容器
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Pt(abs_x * SCALE), Pt(abs_y * SCALE),
                                  Pt(c.w * SCALE), Pt(c.h * SCALE))
    fc = hex_to_rgb(s.get("fillColor", "#f8fafc"))
    sc = hex_to_rgb(s.get("strokeColor", "#cbd5e1"))
    width = float(s.get("strokeWidth", "1") or 1)
    set_fill(body, fc)
    set_line(body, sc, width_pt=width, dashed=False)
    set_no_shadow(body)

    # 标题条 - 解析 raw_value (含 <br>) 后生成多行 run
    title_h = max(start_size, 32)  # 至少 32pt
    tb = slide.shapes.add_textbox(
        Pt(abs_x * SCALE + 14),
        Pt(abs_y * SCALE + 8),
        Pt(c.w * SCALE - 24),
        Pt(title_h - 8),
    )
    set_fill(tb, None)
    set_line(tb, None, width_pt=0)
    set_no_shadow(tb)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    parser = TextParser(c.raw_value,
                        default_size=font_size_of(c),
                        default_bold=True,  # swimlane 标题统一加粗, 与原 drawio 视觉一致
                        default_color=font_color_of(c) or "#0f172a")
    # 单独构造每个 paragraph 以保证 <br> 产生真换行
    lines = parser.get_lines()
    if lines:
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.LEFT
        _runs_to_paragraph(p0, lines[0])
        for ln in lines[1:]:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            _runs_to_paragraph(p, ln)


# ---------------------------------------------------------------------------
# 边 (connector)
# ---------------------------------------------------------------------------

def create_edge(slide, c: Cell, cells: dict[str, Cell], abs_pos: dict[str, tuple[float, float]]) -> None:
    if not (c.source and c.target):
        return
    src = cells.get(c.source)
    tgt = cells.get(c.target)
    if not src or not tgt:
        return

    s = c.style
    stroke_rgb = hex_to_rgb(s.get("strokeColor", "#0f172a"))
    width = float(s.get("strokeWidth", "1") or 1)
    dashed, pat = _line_pattern(s)

    def exit_point(cell: Cell) -> tuple[float, float]:
        ax, ay = abs_pos[cell.id]
        ex = float(cell.style.get("exitX", "0.5") or 0.5)
        ey = float(cell.style.get("exitY", "0.5") or 0.5)
        return (ax + cell.w * ex, ay + cell.h * ey)

    def entry_point(cell: Cell) -> tuple[float, float]:
        ax, ay = abs_pos[cell.id]
        ix = float(cell.style.get("entryX", "0.5") or 0.5)
        iy = float(cell.style.get("entryY", "0.5") or 0.5)
        return (ax + cell.w * ix, ay + cell.h * iy)

    s_pt = exit_point(src)
    t_pt = entry_point(tgt)
    path: list[tuple[float, float]] = [s_pt] + list(c.waypoints) + [t_pt]

    # 删除 waypoints 与端点重合的点
    cleaned: list[tuple[float, float]] = []
    for p in path:
        if not cleaned or (abs(p[0] - cleaned[-1][0]) > 0.1 or abs(p[1] - cleaned[-1][1]) > 0.1):
            cleaned.append(p)

    segments = list(zip(cleaned[:-1], cleaned[1:]))

    last_seg_shape = None
    for i, (p1, p2) in enumerate(segments):
        x1, y1 = p1
        x2, y2 = p2
        dx, dy = x2 - x1, y2 - y1
        if abs(dx) < 0.5 and abs(dy) < 0.5:
            continue
        # python-pptx STRAIGHT 不允许负 width/height.
        # 计算绝对最小包围框,  然后通过显式 begin/end 坐标确定线段方向.
        cx = min(x1, x2) * SCALE
        cy = min(y1, y2) * SCALE
        cw = max(abs(dx) * SCALE, 1)
        ch = max(abs(dy) * SCALE, 1)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          Pt(cx), Pt(cy), Pt(cw), Pt(ch))

        # 找到 cxnSp / xfrm 子元素, 直接覆写 begin/end
        sp = conn._element
        # avLst 是 begin/end 偏移列表 (相对 xfrm left/top)
        # 重新定位 begin 为 (x1 - cx_abs_pt, y1 - cy_abs_pt)
        cnv_sppr = sp.find(qn("p:nvSpPr"))
        # 使用绝对 EMU
        x1_emu = int(x1 * SCALE * 12700)
        y1_emu = int(y1 * SCALE * 12700)
        x2_emu = int(x2 * SCALE * 12700)
        y2_emu = int(y2 * SCALE * 12700)
        # 修改 xfrm.off 和 ext 准确放到 (x1, y1) 至 (x2, y2)
        xfrm = sp.xfrm
        xfrm.off.x = min(x1_emu, x2_emu)
        xfrm.off.y = min(y1_emu, y2_emu)
        xfrm.ext.cx = max(abs(x2_emu - x1_emu), 12700)
        xfrm.ext.cy = max(abs(y2_emu - y1_emu), 12700)
        # 通过 flipH/flipV 让 STRAIGHT begin 落在 (min pt), end 落在 (max pt)
        # 如果 dx<0, begin 在右侧 -> flipH
        if dx < 0:
            xfrm.flipH = True
        if dy < 0:
            xfrm.flipV = True

        set_line(conn, stroke_rgb, width_pt=width, dashed=dashed, dash_pattern=pat)
        last_seg_shape = conn

    if last_seg_shape is not None:
        add_arrow_head_end(last_seg_shape)

    # 标签
    if c.raw_value and c.raw_value.strip():
        mid_i = len(cleaned) // 2
        if mid_i == 0:
            mid_i = 1
        if mid_i >= len(cleaned):
            mid_i = len(cleaned) - 1
        x1, y1 = cleaned[mid_i - 1]
        x2, y2 = cleaned[mid_i]
        midx = (x1 + x2) / 2
        midy = (y1 + y2) / 2
        text = c.raw_value
        # 中文 / 英文 字符宽度估算 - 至少 14pt 每个字符来算标签宽, 避免被裁切
        cn_count = sum(1 for c_ in text if '一' <= c_ <= '鿿')
        en_count = len(text) - cn_count
        approx_w = max(120.0, cn_count * 14 + en_count * 7.5 + 28)
        approx_h = 36
        tb = slide.shapes.add_textbox(
            Pt((midx - approx_w / 2) * SCALE),
            Pt((midy - approx_h / 2) * SCALE),
            Pt(approx_w * SCALE),
            Pt(approx_h * SCALE),
        )
        set_fill(tb, RGBColor(0xff, 0xff, 0xff))
        set_line(tb, None, width_pt=0)
        set_no_shadow(tb)
        tf = tb.text_frame
        tf.margin_left = Pt(2); tf.margin_right = Pt(2)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(12)
        c_obj = hex_to_rgb(c.style.get("fontColor", "#334155") or "#334155")
        if c_obj:
            f.color.rgb = c_obj
        f.bold = False


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------

def main():
    inp, out = resolve_paths()
    if not inp.exists():
        sys.exit(f"找不到输入文件: {inp}")

    cells = load_cells(inp)
    abs_pos = compute_absolute(cells)

    # 计算 bbox - vertex + 边的端点 + waypoints
    min_x = math.inf; min_y = math.inf
    max_x = -math.inf; max_y = -math.inf

    def _consider(x: float, y: float) -> None:
        nonlocal min_x, min_y, max_x, max_y
        if x < min_x: min_x = x
        if y < min_y: min_y = y
        if x > max_x: max_x = x
        if y > max_y: max_y = y

    for cid, c in cells.items():
        if c.is_vertex and not is_edge_label(c) and (c.w > 0.5 or c.h > 0.5):
            ax, ay = abs_pos[cid]
            _consider(ax, ay)
            _consider(ax + c.w, ay + c.h)

    for cid, c in cells.items():
        if not c.is_edge or not c.source or not c.target:
            continue
        src = cells.get(c.source); tgt = cells.get(c.target)
        if not src or not tgt:
            continue
        ax, ay = abs_pos[src.id]
        ex = float(src.style.get("exitX", 0.5) or 0.5); ey = float(src.style.get("exitY", 0.5) or 0.5)
        _consider(ax + src.w * ex, ay + src.h * ey)
        ax, ay = abs_pos[tgt.id]
        ix = float(tgt.style.get("entryX", 0.5) or 0.5); iy = float(tgt.style.get("entryY", 0.5) or 0.5)
        _consider(ax + tgt.w * ix, ay + tgt.h * iy)
        for px, py in c.waypoints:
            _consider(px, py)

    if not (max_x > min_x and max_y > min_y):
        sys.exit("无法从图中获得有效 bbox")

    # 平移到画布原点
    offset_x = -min_x + 20
    offset_y = -min_y + 20

    # 所有 shape 经过 offset 后, 总宽度 = max_x - min_x + 40
    slide_w_pt = (max_x - min_x + 40)
    slide_h_pt = (max_y - min_y + 40)

    prs = Presentation()
    prs.slide_width = Emu(int(slide_w_pt * 12700))
    prs.slide_height = Emu(int(slide_h_pt * 12700))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)

    # 应用 offset 到 abs_pos
    abs_pos2 = {k: (x + offset_x, y + offset_y) for k, (x, y) in abs_pos.items()}
    waypoints2: dict[str, list[tuple[float, float]]] = {}
    for cid, c in cells.items():
        if c.waypoints:
            waypoints2[cid] = [(px + offset_x, py + offset_y) for px, py in c.waypoints]

    # 先画所有 vertex, 后画 edge
    for cid, c in cells.items():
        if not c.is_vertex or is_edge_label(c):
            continue
        if c.w < 0.5 and c.h < 0.5:
            continue
        ax, ay = abs_pos2[cid]
        if shape_kind(c) == "swimlane":
            create_swimlane(slide, c, ax, ay)
        else:
            create_vertex(slide, c, ax, ay)

    for cid, c in cells.items():
        if c.is_edge and c.source and c.target:
            # 临时替换 waypoints
            wp_orig = c.waypoints
            c.waypoints = waypoints2.get(cid, wp_orig)
            create_edge(slide, c, cells, abs_pos2)
            c.waypoints = wp_orig

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"生成完成: {out}")
    print(f"  幻灯片大小: {slide_w_pt / 72:.2f}\" x {slide_h_pt / 72:.2f}\"")
    print(f"  节点数: {sum(1 for c in cells.values() if c.is_vertex and not is_edge_label(c) and (c.w > 0.5 or c.h > 0.5))}")
    print(f"  边数: {sum(1 for c in cells.values() if c.is_edge)}")


if __name__ == "__main__":
    main()
